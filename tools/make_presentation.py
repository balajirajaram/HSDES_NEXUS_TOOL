"""Generate presentation assets for the Auto HSD Analyser:
  1) a 16:9 architecture / data-flow PNG, and
  2) a PowerPoint deck that embeds it.
Run: python tools/make_presentation.py
"""
import os

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(HERE, "presentation")
os.makedirs(OUT, exist_ok=True)
PNG = os.path.join(OUT, "auto_hsd_analyser_flow.png")
SHOT = os.path.join(OUT, "report_screenshot.png")
PPTX = os.path.join(OUT, "Auto_HSD_Analyser.pptx")

# Intel-ish palette
INK = "#1b2a4a"
BLUE = "#0068b5"
LBLUE = "#d7ebfb"
TEAL = "#00857d"
LTEAL = "#d3efec"
AMBER = "#b26a00"
LAMBER = "#fbeacd"
GREEN = "#2e7d32"
GREY = "#5b6b7f"
LGREY = "#eef2f7"


def _box(ax, x, y, w, h, title, lines, edge, fill, tsize=12, lsize=9.2):
    ax.add_patch(FancyBboxPatch(
        (x, y), w, h, boxstyle="round,pad=0.006,rounding_size=0.014",
        linewidth=2, edgecolor=edge, facecolor=fill, zorder=2))
    ax.text(x + w / 2, y + h - 0.045, title, ha="center", va="top",
            fontsize=tsize, fontweight="bold", color=INK, zorder=3)
    if lines:
        ax.text(x + w / 2, y + h - 0.105, "\n".join(lines), ha="center", va="top",
                fontsize=lsize, color="#26364f", zorder=3, linespacing=1.5)


def _arrow(ax, p0, p1, color=BLUE, style="-|>", lw=2.2, rad=0.0):
    ax.add_patch(FancyArrowPatch(
        p0, p1, arrowstyle=style, mutation_scale=16, linewidth=lw,
        color=color, connectionstyle=f"arc3,rad={rad}", zorder=1))


def build_png():
    fig, ax = plt.subplots(figsize=(16, 9), dpi=150)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    fig.patch.set_facecolor("white")

    # Header band
    ax.add_patch(FancyBboxPatch((0.0, 0.905), 1.0, 0.095,
                 boxstyle="square,pad=0", linewidth=0, facecolor=INK, zorder=0))
    ax.text(0.5, 0.963, "Auto HSD Analyser — How It Works",
            ha="center", va="center", fontsize=22, fontweight="bold", color="white")
    ax.text(0.5, 0.925, "Give it an HSD ID → it reads, decodes, correlates and reports an "
            "engineer-grade root-cause triage — then learns.",
            ha="center", va="center", fontsize=11.5, color="#c9d6ea")

    # 1. Input
    _box(ax, 0.03, 0.585, 0.155, 0.20, "1 · Input",
         ["HSD ID", "(+ optional", "symptom / logs)"], BLUE, LBLUE, tsize=13)

    # 2. Auth + fetch
    _box(ax, 0.225, 0.545, 0.20, 0.28, "2 · Secure Fetch",
         ["Kerberos SSO — no", "password / token", "", "HSDES: title, desc,",
          "full comment thread,", "attachments"], BLUE, LBLUE, tsize=13)

    # 3. Four parallel analysis engines
    ex, ew, eh, step = 0.475, 0.235, 0.12, 0.14
    tops = [0.75, 0.61, 0.47, 0.33]
    _box(ax, ex, tops[0], ew, eh, "3a · Comment Investigation",
         ["Rebuilds narrative → converged root cause"],
         TEAL, LTEAL, tsize=12, lsize=8.8)
    _box(ax, ex, tops[1], ew, eh, "3b · Log & Attachment Decode",
         ["SOL / PythonSV / POST → Intel decoder DBs"],
         TEAL, LTEAL, tsize=12, lsize=8.8)
    _box(ax, ex, tops[2], ew, eh, "3c · Axon Recordings",
         ["Linked recordings → SVTools signatures"],
         TEAL, LTEAL, tsize=12, lsize=8.8)
    _box(ax, ex, tops[3], ew, eh, "3d · Knowledge Base",
         ["KB recall + similar past HSDs & fixes"],
         TEAL, LTEAL, tsize=12, lsize=8.8)
    ectr = [t + eh / 2 for t in tops]

    # 4. Reason / correlate
    _box(ax, 0.745, 0.50, 0.205, 0.30, "4 · Correlate & Reason",
         ["Signal vs. noise ranking", "", "LLM prose  —or—", "deterministic OFFLINE",
          "mode (no LLM needed)"], AMBER, LAMBER, tsize=13)

    # 5. Report
    _box(ax, 0.745, 0.205, 0.205, 0.22, "5 · Triage Report",
         ["Markdown + HTML", "Debug summary · timeline ·", "root cause · next actions"],
         GREEN, "#dcefe0", tsize=13)

    # 6. Learn feedback
    _box(ax, 0.475, 0.095, 0.235, 0.10, "6 · Learn",
         ["Every ticket written back to the KB — smarter over time"],
         GREY, LGREY, tsize=12, lsize=8.4)

    # arrows
    _arrow(ax, (0.185, 0.685), (0.225, 0.685))
    for cy, rad in zip(ectr, (-0.18, -0.05, 0.06, 0.20)):
        _arrow(ax, (0.425, 0.685), (0.475, cy), rad=rad)
    for cy in ectr:
        _arrow(ax, (0.710, cy), (0.745, 0.65), color=AMBER,
               rad=0.05 if cy > 0.65 else -0.05)
    _arrow(ax, (0.8475, 0.50), (0.8475, 0.425), color=GREEN)
    # learn loop: report -> learn -> KB
    _arrow(ax, (0.745, 0.255), (0.71, 0.16), color=GREY, rad=0.1)
    _arrow(ax, (0.475, 0.16), (0.44, 0.39), color=GREY, rad=0.25, style="-|>")

    # footer legend
    ax.text(0.03, 0.02, "Blue = data intake   ·   Teal = parallel analysis engines   ·   "
            "Amber = reasoning   ·   Green = output   ·   Grey = continuous learning",
            ha="left", va="center", fontsize=10, color=GREY)

    # footer legend
    ax.text(0.03, 0.02, "Blue = data intake   ·   Teal = parallel analysis engines   ·   "
            "Amber = reasoning   ·   Green = output   ·   Grey = continuous learning",
            ha="left", va="center", fontsize=10, color=GREY)

    fig.savefig(PNG, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("PNG :", PNG)


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY = RGBColor(0x1B, 0x2A, 0x4A)
    BLU = RGBColor(0x00, 0x68, 0xB5)
    GRY = RGBColor(0x44, 0x55, 0x66)
    WHT = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def band(slide, title, sub=None):
        bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        tf = bar.text_frame; tf.margin_left = Inches(0.5); tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHT
        if sub:
            sp = tf.add_paragraph(); sp.text = sub
            sp.font.size = Pt(13); sp.font.color.rgb = RGBColor(0xC9, 0xD6, 0xEA)

    def bullets(slide, items, top=1.5, left=0.7, width=12.0, size=18, height=5.4):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            lvl = 0
            if isinstance(it, tuple):
                it, lvl = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = it; p.level = lvl
            p.font.size = Pt(size - 3 * lvl)
            p.font.color.rgb = GRY if lvl else NAVY
            p.space_after = Pt(7)
        return tb

    # ---- 1 Title ----
    s = add_slide()
    rect = s.shapes.add_shape(1, 0, 0, SW, SH)
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.6))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Auto HSD Analyser"
    p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = WHT
    p = tf.add_paragraph(); p.text = "Automated, engineer-grade HSD triage for Intel server platforms"
    p.font.size = Pt(22); p.font.color.rgb = RGBColor(0x9F, 0xC3, 0xE8)
    p = tf.add_paragraph(); p.text = "Read · Decode · Correlate · Report · Learn"
    p.font.size = Pt(16); p.font.color.rgb = RGBColor(0xC9, 0xD6, 0xEA)
    p = tf.add_paragraph(); p.text = "\nPresented by: ______________          Date: ____________"
    p.font.size = Pt(14); p.font.color.rgb = RGBColor(0xC9, 0xD6, 0xEA)

    # ---- 2 Problem ----
    s = add_slide(); band(s, "The Problem", "HSD triage is slow, manual and inconsistent")
    bullets(s, [
        "Every HSD means reading the whole comment thread, opening SOL / PythonSV / POST logs, and decoding raw error codes by hand.",
        "Decisive evidence is buried — register values, repro steps, experiment matrices, the converged root cause.",
        "Verbose serial logs are full of incidental noise (WHEA / MCA / PCIe) that hides the real failure.",
        "Past-ticket knowledge lives in people's heads, not in a searchable, reusable form.",
        "Net: hours per ticket, uneven quality, and re-investigation of already-solved issues.",
    ], size=19)

    # ---- 3 What it does ----
    s = add_slide(); band(s, "What It Does", "One HSD ID in → a structured root-cause triage out")
    bullets(s, [
        "READ — pulls the full ticket over Kerberos SSO: title, description, entire comment thread, attachments.",
        "DECODE — scans SOL / PythonSV / POST logs; decodes MCE/MCA (MCACOD · MSCOD), EWL, RC-Fatal, MCHECK, POST; pulls Axon SVTools signatures.",
        "CORRELATE — reconstructs the investigation, converges on the root cause, separates the real failure from log noise.",
        "REPORT — clean Markdown + HTML triage: debug summary, timeline, findings, root cause, next actions.",
        "LEARN — writes every ticket into a self-learning Knowledge Base; recall improves over time.",
    ], size=19)

    # ---- 4 How it works (diagram) ----
    s = add_slide(); band(s, "How It Works", "End-to-end data flow")
    s.shapes.add_picture(PNG, Inches(0.35), Inches(1.35), width=Inches(12.6))

    def panel(slide, x, w, title, tcolor, items, top=1.7, h=4.7):
        box = slide.shapes.add_shape(5, Inches(x), Inches(top), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFB)
        box.line.color.rgb = tcolor; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.22)
        p = tf.paragraphs[0]; p.text = title
        p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = tcolor
        p.space_after = Pt(6)
        for it in items:
            pp = tf.add_paragraph(); pp.text = "•  " + it
            pp.font.size = Pt(14.5); pp.font.color.rgb = GRY; pp.space_before = Pt(7)

    # ---- 5 Two paths to root cause ----
    s = add_slide(); band(s, "How It Finds the Root Cause",
                          "Two paths — it always returns an answer")
    panel(s, 0.55, 5.9, "A · When engineers have investigated", BLU, [
        "Reconstructs the comment thread in order.",
        "Reports the team's converged root cause, verbatim.",
        "Cites who said it + how proven (hypothesis vs confirmed).",
        "Surfaces the decisive register / experiment evidence.",
    ])
    panel(s, 6.9, 5.9, "B · When there are no comments", RGBColor(0x00, 0x85, 0x7D), [
        "Analyses the attached logs end-to-end.",
        "Decodes MCE / MCA codes → bank, IP, error class.",
        "Folds in Axon SVTools failure signatures.",
        "Proposes an evidence-based hypothesis + data to confirm.",
    ])

    # ---- 6 Key capabilities ----
    s = add_slide(); band(s, "Key Capabilities", "Reads like a human debug engineer")
    bullets(s, [
        "Comment-thread analyzer — rebuilds who observed / tried / concluded what, converges on the cause + disposition.",
        "Deep log decode — MCE/MCA (MCACOD · MSCOD → bank / IP / class), EWL, RC-Fatal, MCHECK, POST checkpoints, boot/stage progress.",
        "Axon integration — pulls linked recordings and their SVTools failure signatures as precise corroborating evidence.",
        "Signal-vs-noise ranking — flags incidental WHEA / MCA / PCIe telemetry so the real failure leads.",
        "Deterministic OFFLINE mode — full report with zero LLM dependency; optional LLM adds prose.",
        "Self-learning KB — recalls similar past HSDs with their root cause and fix.",
    ], size=17)

    # ---- 7 Sample report anatomy ----
    s = add_slide(); band(s, "Anatomy of a Report", "What your manager will see")
    bullets(s, [
        "Debug Summary — RESULT, boot/stage progress, metrics, hypothesis, conclusion, next actions.",
        "Investigation Summary — compact timeline, key evidence, converged root cause, disposition.",
        "Findings — failure signatures (noise flagged), decoded MCA / BIOS / POST, confidence score.",
        "Root Cause + Validation — verdict, provenance, and what's needed to confirm it.",
        "Next Actions — targeted checks + PythonSV / rdmsr commands to close it out.",
        "Appendix — full narrative, decoded log evidence, KB matches, similar HSDs.",
    ])

    # ---- 8 Demo ----
    s = add_slide(); band(s, "Live Demo", "Sample: HSD 16031066261 — Warm Reset / GBLRST_CAUSE1")
    bullets(s, [
        "Root cause: S3M version mismatch (Base/Fit-loaded UP vs OSPL UP).",
        "Evidence: GBLRST_CAUSE1 = 0x4 (pass) vs 0x0 (fail) + 6-row experiment matrix.",
        "Noise flag: WHEA / MCA / PCIe marked incidental.",
        "Disposition: root-caused; handed to S3M team; ticket open.",
    ], left=0.55, width=5.4, size=15)
    if os.path.exists(SHOT):
        s.shapes.add_picture(SHOT, Inches(6.15), Inches(1.45), width=Inches(6.7))

    # ---- 9 Benefits ----
    s = add_slide(); band(s, "Benefits & Impact",
                          "≈ 2 min per ticket vs 2–4 hrs manual (illustrative estimate — tune to your data)")
    bullets(s, [
        "Faster — minutes vs hours per ticket; reading, decoding and correlation are automated.",
        "Consistent — engineer-grade output every time, independent of who runs it.",
        "Complete — decisive register / MCA / experiment evidence surfaced automatically, nothing missed.",
        "Always answers — cites the team's cause, or proposes one from logs + MCE/MCA + Axon.",
        "Secure — Kerberos SSO, runs locally on any Intel-domain Windows machine.",
        "Compounding — the Knowledge Base makes every future triage faster and smarter.",
    ], size=17)

    # ---- 10 Data sources & provenance ----
    s = add_slide(); band(s, "Data Sources & Provenance",
                          "Every fact in the report is traceable — nothing leaves the Intel network")
    bullets(s, [
        "Root cause & investigation — the HSDES ticket itself: description + full comment thread, read live over Kerberos SSO and cited per author.",
        "Attached logs — SOL / PythonSV / POST logs from the HSDES attachments API, plus linked Axon recordings (SVTools failure signatures).",
        "Silicon MCE / MCA decode — bundled Intel decoder databases in app/decoders/, derived from Intel MCA / RAS customer documentation:",
        ("MCA 712 codes · MCHECK 496 · EWL 124 · RC-Fatal 50 · IPSD 28 · POST checkpoints (MCACOD / MSCOD → bank, IP, error class, action).", 1),
        "Similar past issues & fixes — HSDES search + a local self-learning Knowledge Base (kb/).",
        "All processing is local; the report labels each cause with its source and how proven it is.",
    ], size=16)

    # ---- 9 Roadmap ----
    s = add_slide(); band(s, "Roadmap")
    bullets(s, [
        "Broaden product coverage (DMR / COR) and seed the KB from more master queries.",
        "Deeper Axon / PythonSV automation and richer decoder databases.",
        "Optional LLM reasoning layer enabled by default where an endpoint is available.",
        "One-click ticket update: write the converged root cause back to the HSDES field.",
        "Team rollout + feedback loop to keep improving accuracy.",
    ])

    def qa(slide, pairs, top=1.45, size=15):
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(5.6))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for q, a in pairs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = "Q:  " + q
            p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = BLU
            p.space_before = Pt(6); p.space_after = Pt(2)
            pa = tf.add_paragraph(); pa.text = "A:  " + a
            pa.font.size = Pt(size - 1); pa.font.color.rgb = GRY
            pa.space_after = Pt(6)

    # ---- 10 Anticipated Q&A (part 1) ----
    s = add_slide(); band(s, "Anticipated Questions (1 of 2)", "Prepared answers for the demo")
    qa(s, [
        ("Does any Intel data leave the network / go to an external AI?",
         "No. It runs locally on your Intel-domain machine and talks only to HSDES. LLM reasoning is optional and off by default — the demo runs in deterministic OFFLINE mode."),
        ("How secure is authentication?",
         "Kerberos SSO — it authenticates as you, with no password and no stored token. Any teammate on the Intel domain can run it as themselves."),
        ("Can we trust the root cause it reports?",
         "Yes — it doesn't guess. It repeats the root cause the engineers actually wrote in the ticket, cites who said it, and labels how proven it is (e.g. 'unvalidated hypothesis' vs 'confirmed'). If no cause is stated, it says so instead of inventing one."),
        ("Does it modify or close the HSD ticket?",
         "No — it is read-only today. Writing the root cause back to the HSDES field is on the roadmap as an explicit, opt-in action."),
    ])

    # ---- 11 Anticipated Q&A (part 2) ----
    s = add_slide(); band(s, "Anticipated Questions (2 of 2)", "Prepared answers for the demo")
    qa(s, [
        ("How does it separate the real failure from log noise?",
         "When the comment thread has converged on a cause, broad signatures (WHEA / MCA / PCIe) are explicitly flagged as incidental background telemetry, so the real failure leads."),
        ("What if the comments are thin or missing?",
         "It falls back to decoded log evidence, boot/stage progress, MCA/BIOS/POST decode, and Knowledge-Base recall of similar past HSDs — you still get a structured triage."),
        ("How long does a ticket take, and what does it cost?",
         "Seconds to a couple of minutes per ticket versus hours by hand. It uses only open-source Python libraries and runs on your existing machine — no new licensing."),
        ("Which platforms and domains are covered?",
         "Server platforms today (GNR / SRF / CWF) across RAS/MCA, UPI, memory, PCIe/CXL, power, BIOS/IFWI/BMC, OS/driver; DMR / COR are ready to add."),
        ("How does the 'learning' work?",
         "Every analysed ticket is written into a local self-learning Knowledge Base, so recall of similar issues and their fixes improves over time."),
    ], size=14)

    # ---- 12 Thank you ----
    s = add_slide()
    rect = s.shapes.add_shape(1, 0, 0, SW, SH)
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.6))
    tf = t.text_frame
    p = tf.paragraphs[0]; p.text = "Thank you"
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHT
    p = tf.add_paragraph(); p.text = "Questions & discussion"
    p.font.size = Pt(20); p.font.color.rgb = RGBColor(0x9F, 0xC3, 0xE8)

    prs.save(PPTX)
    print("PPTX:", PPTX)


if __name__ == "__main__":
    build_png()
    build_pptx()
